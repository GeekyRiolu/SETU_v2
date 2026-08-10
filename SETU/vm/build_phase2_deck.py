#!/usr/bin/env python3
"""Update SETU.pptx (slides Drawbacks -> Conclusion) to match the system we
actually built: SeqKD distillation + English pivot, 11/22 languages in Phase 1,
real interfaces and metrics. Preserves theme, tables and images; writes a new
file so an open copy is never clobbered.

    .venv/bin/python vm/build_phase2_deck.py
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

SRC = "vm/SETU.pptx"
DST = "vm/SETU_updated.pptx"
FONT = "Times New Roman"

# ---- body-text rewrites, keyed by slide index (bullets, 6-8 lines) ----
BODY = {
    11: [  # Proposed System
        "SETU distils a strong teacher (IndicTrans2) into compact 52M-parameter student translators built for offline use.",
        "The core method is Sequence-Level Knowledge Distillation (SeqKD): the student learns the teacher's own 1-best translations, not noisy web references.",
        "Our experiments show SeqKD decisively outperforms preference/DPO distillation for compact Indic machine translation.",
        "Indic-to-Indic translation is served by pivoting through English (source -> English -> target).",
        "Students are quantised to INT4 ONNX (about 104 MB per direction) for on-device, fully offline inference.",
        "Phase 1 delivers 11 of the 22 scheduled languages, both directions, through a REST API, an offline PWA, a web app and a CLI.",
    ],
    13: [  # Conceptual/Analysis Modelling - OO models
        "Object-Oriented models describe SETU's structure and behaviour.",
        "Use Case Diagram: actors are End User, ML Engineer and Edge Device; use cases are Translate Text, Distil Teacher Targets, Train Student and Quantise Model.",
        "Sequence Diagram: source text -> tokenisation -> student ONNX inference -> detokenised translation, with an English pivot for Indic-Indic pairs.",
        "Activity Diagram: corpus loading, teacher distillation and student SeqKD training run offline, followed by quantisation and packaging.",
        "Class Diagram: key classes are InferenceEngine, ONNXTranslator, IndicTrans2Teacher, Distiller, StudentTokenizer, ModelConfig and TranslationResult.",
        "At run time the user touches only the Inference Engine; the training classes are used offline by the ML Engineer.",
    ],
    14: [  # Conceptual/Analysis Modelling - structured models
        "Structured development models describe SETU's data and state flow.",
        "Data Flow Diagram: Samanantar corpus -> teacher 1-best distillation -> SeqKD student training -> INT4 ONNX quantisation -> offline inference.",
        "State Chart: a model moves Teacher-distilled -> SeqKD-trained -> Quantised (INT4 ONNX) -> Deployed on device.",
        "ER Model: entities are Language, SentencePair, ModelVersion and TranslationResult, linked by training and evaluation relationships.",
        "The English pivot adds a routing state for Indic-Indic pairs that have no direct model.",
        "Evaluation compares student vs teacher BLEU/chrF and gates deployment on the 80% ratio target.",
    ],
    19: [  # Module Decomposition
        "Module 1 - Corpus Loader: ingests Samanantar parallel data, normalises scripts and aligns sentence pairs.",
        "Module 2 - Teacher Distiller: runs IndicTrans2 (dist-200M) to generate 1-best translations as SeqKD targets.",
        "Module 3 - SeqKD Trainer: trains the 52M student on teacher targets with PAD masking, LR warmup and gradient clipping.",
        "Module 4 - Quantiser: applies INT8/INT4 post-training quantisation and exports to ONNX (about 104 MB).",
        "Module 5 - Inference Engine: ONNX Runtime plus SentencePiece; routes Indic-Indic pairs through English.",
        "Module 6 - REST API: a FastAPI service exposing /translate, /languages, /models and /health.",
        "Module 7 - Frontends: an offline PWA, a Next.js web app and a CLI, all sharing one engine.",
        "Module 8 - Evaluator: sacrebleu BLEU/chrF and the student/teacher ratio on held-out and FLORES sets.",
    ],
    22: [  # Algorithm Design
        "Step 1 - Teacher Distillation: IndicTrans2 greedily generates 1-best translations over the parallel corpus.",
        "Step 2 - SeqKD Training: the student is trained with token-level cross-entropy on the teacher's outputs (Kim & Rush, 2016), masking PAD tokens.",
        "Step 3 - Stability: LR warmup, label smoothing and gradient clipping prevent the collapse seen at higher learning rates.",
        "Step 4 - Key Result: SeqKD beats reference/DPO training; the student/teacher ratio scales 0.607 -> 0.764 -> 0.806 as data grows 100k -> 250k -> 500k.",
        "Step 5 - Quantisation: post-training INT8/INT4 export to ONNX (about 104 MB, roughly 4x compression), fully offline.",
        "Step 6 - English Pivot: Indic-to-Indic translation chains two students, source -> English -> target.",
        "Step 7 - Evaluation: BLEU/chrF with the 80% student/teacher ratio as the quality gate.",
    ],
    24: [  # Conclusion
        "SETU delivers offline, on-device translation across the scheduled languages of India by distilling IndicTrans2 into compact 52M students.",
        "Its core finding: sequence-level KD decisively outperforms preference/DPO distillation for compact Indic MT; reference noise, not the method, was the bottleneck.",
        "Phase 1 covers 11 of 22 languages, both directions; Indic-to-Indic works via English pivot.",
        "Students are INT4 ONNX (about 104 MB), translate in under 500 ms on CPU, and run with zero network calls.",
        "12 of the first 18 directions already meet the 80% teacher-BLEU target; the rest are near-miss.",
        "The system is delivered as a REST API, an offline PWA, a Next.js web app and a CLI.",
        "Phase 2 extends to the remaining 11 languages and adds FLORES/COMET benchmarking.",
    ],
}

RESULTS = [  # new slide inserted before Conclusion
    "11 of 22 scheduled languages are covered in Phase 1 (Assamese and Punjabi in progress); each direction is its own 52M student.",
    "12 of the first 18 directions meet the 80% student/teacher BLEU ratio; the strongest is Marathi->English at 0.955.",
    "Every model is INT4 ONNX at about 104 MB and translates in roughly 40-320 ms on commodity CPU, fully offline.",
    "English pivot makes all Indic-to-Indic pairs work (for example Hindi->Bengali via English), not just Indic-English.",
    "Two training collapses (Telugu->English, Malayalam->English) were diagnosed and fixed by lowering the SFT learning rate to 3e-4.",
    "Working demo: a Next.js web app and offline PWA calling a FastAPI engine that auto-discovers every trained model.",
]

# ---- targeted table-cell replacements: {slide_idx: {match_substring: new_cell_text}} ----
TABLES = {
    12: {
        "NVIDIA A100": "Training GPU: NVIDIA A40 (16 GB vGPU slice)",
        "DPO training library": "optimum + ONNX Runtime (INT8/INT4 export + inference)",
        "IndicTrans2 model weights": "IndicTrans2 dist-200M weights (ungated teacher)",
        "DPO training library (TRL / custom implementation)": "optimum + ONNX Runtime (INT8/INT4 export + inference)",
    },
    16: {
        "Translate text across all 22": "Translate across Indian languages offline (11 of 22 in Phase 1, both directions)",
        "Model size ≤ 200 MB after quantization for edge deployment": "Model size <= 200 MB after quantization (about 104 MB achieved, INT4)",
        "Provide confidence scores": "Support Indic-to-Indic translation via English pivot",
    },
    20: {
        "python setup.py --src": 'python setu_cli.py --src hi --tgt en --text "..."',
        "Dropdown for languages": "Next.js app: language pickers, swap, live latency + on-device badge",
        "Route:/translate": "Method: POST  Route: /translate\nAlso: GET /languages, /models, /health\nPayload: {source_lang, target_lang, text}",
    },
}


def style_body(tf, bullets, size=24):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink to fit, no overflow
    except Exception:
        pass
    tf.clear()
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.level = 0
        para.line_spacing = 1.5
        para.alignment = PP_ALIGN.JUSTIFY
        for run in para.runs:
            run.font.name = FONT
            run.font.size = Pt(size)


def find_body(slide):
    # idx 0 is the title placeholder; the body is the next placeholder.
    # (Do not compare against slide.shapes.title by identity: python-pptx
    #  returns distinct wrapper objects for the same placeholder.)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame:
            return ph
    return None


def set_cell(cell, text):
    cell.text = text
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = FONT


prs = Presentation(SRC)

for idx, bullets in BODY.items():
    body = find_body(prs.slides[idx])
    if body is None:
        print(f"WARN: no body on slide {idx}")
        continue
    style_body(body.text_frame, bullets)
    print(f"slide {idx}: rewrote body ({len(bullets)} bullets)")

for idx, repls in TABLES.items():
    for sh in prs.slides[idx].shapes:
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    for match, new in repls.items():
                        if match in cell.text:
                            set_cell(cell, new)
                            print(f"slide {idx}: table cell updated ({match[:28]}...)")
                            break

# ---- add the Results slide, then move it to sit just before Conclusion ----
concl_idx = 24  # "Conclusion"
layout = next(l for l in prs.slide_layouts if l.name == "OBJECT")
slide = prs.slides.add_slide(layout)
if slide.shapes.title is not None:
    slide.shapes.title.text = "Implementation & Results (Phase 1)"
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.name = FONT
body = find_body(slide)
if body is not None:
    style_body(body.text_frame, RESULTS)

xml_slides = prs.slides._sldIdLst
entries = list(xml_slides)
new_entry = entries[-1]
xml_slides.remove(new_entry)
xml_slides.insert(concl_idx, new_entry)  # results before Conclusion
print(f"inserted Results slide at index {concl_idx} (before Conclusion)")

prs.save(DST)
print(f"\nSaved -> {DST}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slide entries)")
